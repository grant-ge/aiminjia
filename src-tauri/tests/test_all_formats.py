#!/usr/bin/env python3
"""
Comprehensive test suite for all document formats: Word, Excel, PPT, PDF.
Tests parsing logic for edge cases and validates error handling.
"""

import json
import sys
from pathlib import Path

# Check dependencies
missing = []
try:
    from docx import Document
except ImportError:
    missing.append("python-docx")

try:
    from openpyxl import Workbook
except ImportError:
    missing.append("openpyxl")

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    missing.append("python-pptx")

try:
    import pdfplumber
except ImportError:
    missing.append("pdfplumber")

try:
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas
    from reportlab.platypus import Table, TableStyle
except ImportError:
    missing.append("reportlab")

if missing:
    print(f"ERROR: Missing dependencies: {', '.join(missing)}")
    print(f"Install with: pip install {' '.join(missing)}")
    sys.exit(1)

import pandas as pd

# Test output directory
TEST_DIR = Path(__file__).parent / "fixtures"
TEST_DIR.mkdir(parents=True, exist_ok=True)


# ============================================================
# Excel Test Cases
# ============================================================

def create_excel_merged_cells():
    """Excel: Merged cells in header"""
    wb = Workbook()
    ws = wb.active
    ws.title = "MergedHeader"

    # Merge A1:E1
    ws.merge_cells('A1:E1')
    ws['A1'] = 'Merged Header'

    # Data rows
    for i in range(2, 5):
        for j in range(1, 6):
            ws.cell(row=i, column=j, value=f'R{i}C{j}')

    path = TEST_DIR / "excel_merged.xlsx"
    wb.save(path)
    print(f"✓ Created: {path}")
    return path


def create_excel_multiple_sheets():
    """Excel: Multiple sheets with different structures"""
    wb = Workbook()

    # Sheet 1: 3 columns
    ws1 = wb.active
    ws1.title = "Employees"
    ws1.append(['Name', 'Age', 'Dept'])
    ws1.append(['Alice', 30, 'Eng'])
    ws1.append(['Bob', 25, 'Sales'])

    # Sheet 2: 5 columns
    ws2 = wb.create_sheet("Sales")
    ws2.append(['Product', 'Q1', 'Q2', 'Q3', 'Q4'])
    ws2.append(['Widget', 100, 200, 150, 300])

    path = TEST_DIR / "excel_multi_sheet.xlsx"
    wb.save(path)
    print(f"✓ Created: {path}")
    return path


def create_excel_empty_rows():
    """Excel: Empty rows and columns"""
    wb = Workbook()
    ws = wb.active

    ws.append(['Name', 'Value', 'Notes'])
    ws.append(['Item1', 100, None])
    ws.append([None, None, None])  # Empty row
    ws.append(['Item2', None, 'Note'])

    path = TEST_DIR / "excel_empty.xlsx"
    wb.save(path)
    print(f"✓ Created: {path}")
    return path


def create_excel_large():
    """Excel: Large file (10k rows) to test performance"""
    wb = Workbook()
    ws = wb.active

    ws.append(['ID', 'Name', 'Value', 'Category'])
    for i in range(1, 10001):
        ws.append([i, f'Item{i}', i * 10, f'Cat{i % 5}'])

    path = TEST_DIR / "excel_large.xlsx"
    wb.save(path)
    print(f"✓ Created: {path}")
    return path


# ============================================================
# PPT Test Cases
# ============================================================

def create_ppt_with_tables():
    """PPT: Slides with tables"""
    prs = Presentation()

    # Slide 1: Table
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    title = slide1.shapes.title
    title.text = "Sales Data"

    rows, cols = 3, 4
    left = top = Inches(1)
    width = Inches(8)
    height = Inches(2)

    table = slide1.shapes.add_table(rows, cols, left, top, width, height).table
    table.cell(0, 0).text = 'Product'
    table.cell(0, 1).text = 'Q1'
    table.cell(0, 2).text = 'Q2'
    table.cell(0, 3).text = 'Q3'
    table.cell(1, 0).text = 'Widget'
    table.cell(1, 1).text = '100'
    table.cell(1, 2).text = '200'
    table.cell(1, 3).text = '150'

    # Slide 2: Text only
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Summary"
    slide2.placeholders[1].text = "This is a text-only slide."

    path = TEST_DIR / "ppt_tables.pptx"
    prs.save(path)
    print(f"✓ Created: {path}")
    return path


def create_ppt_text_only():
    """PPT: Text-only slides"""
    prs = Presentation()

    for i in range(3):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Slide {i+1}"
        slide.placeholders[1].text = f"Content for slide {i+1}"

    path = TEST_DIR / "ppt_text.pptx"
    prs.save(path)
    print(f"✓ Created: {path}")
    return path


# ============================================================
# PDF Test Cases
# ============================================================

def create_pdf_with_table():
    """PDF: Document with table (using reportlab)"""
    from reportlab.lib import colors

    path = TEST_DIR / "pdf_table.pdf"
    c = canvas.Canvas(str(path), pagesize=letter)

    # Title
    c.setFont("Helvetica-Bold", 16)
    c.drawString(100, 750, "Sales Report")

    # Table data
    data = [
        ['Product', 'Q1', 'Q2', 'Q3'],
        ['Widget', '100', '200', '150'],
        ['Gadget', '80', '120', '90'],
    ]

    # Create table
    table = Table(data)
    table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 12),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
        ('GRID', (0, 0), (-1, -1), 1, colors.black),
    ]))

    # Draw table
    table.wrapOn(c, 400, 200)
    table.drawOn(c, 100, 600)

    c.save()
    print(f"✓ Created: {path}")
    return path


def create_pdf_text_only():
    """PDF: Text-only document"""
    path = TEST_DIR / "pdf_text.pdf"
    c = canvas.Canvas(str(path), pagesize=letter)

    c.setFont("Helvetica", 12)
    y = 750
    for i in range(10):
        c.drawString(100, y, f"This is paragraph {i+1} of the text-only PDF.")
        y -= 20

    c.save()
    print(f"✓ Created: {path}")
    return path


# ============================================================
# Test Execution
# ============================================================

def test_excel_parsing(xlsx_path):
    """Test Excel parsing logic (mimics parser.rs)"""
    print(f"\n--- Testing: {xlsx_path.name} ---")

    try:
        from openpyxl import load_workbook

        wb = load_workbook(xlsx_path, read_only=True, data_only=True)
        ws = wb.active
        actual_row_count = ws.max_row

        rows = list(ws.iter_rows(max_row=10001, values_only=True))
        wb.close()

        if rows:
            header = [str(c).strip() if c is not None else f'col_{i}' for i, c in enumerate(rows[0])]
            df = pd.DataFrame(rows[1:10001], columns=header)
            actual_data_rows = max(0, actual_row_count - 1)

            print(f"  ✓ Parsed: {actual_data_rows} rows, {len(header)} columns")
            print(f"  Header: {header[:5]}{'...' if len(header) > 5 else ''}")
            print(f"  Sample: {df.head(2).to_dict('records')}")
            return {"type": "table", "rows": actual_data_rows, "cols": len(header)}
        else:
            print(f"  ✓ Empty file")
            return {"type": "empty"}

    except Exception as e:
        print(f"  ✗ FAILED: {e}")
        import traceback
        traceback.print_exc()
        return {"type": "error", "error": str(e)}


def test_ppt_parsing(pptx_path):
    """Test PPT parsing logic (mimics parser.rs)"""
    print(f"\n--- Testing: {pptx_path.name} ---")

    try:
        from pptx import Presentation

        prs = Presentation(pptx_path)
        total_slides = len(prs.slides)

        # Extract tables
        best_table = None
        best_row_count = 0

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    if len(table.rows) > best_row_count:
                        best_row_count = len(table.rows)
                        best_table = table

        if best_table:
            all_rows = []
            header = None

            for i, row in enumerate(best_table.rows):
                cells = [cell.text.strip() for cell in row.cells]
                if header is None and any(cells):
                    header = cells
                elif header is not None:
                    all_rows.append(cells)

            print(f"  ✓ Extracted table: {len(all_rows)} rows, {len(header)} columns from {total_slides} slides")
            print(f"  Header: {header}")
            return {"type": "table", "rows": len(all_rows), "cols": len(header)}
        else:
            # Extract text
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                text_parts.append(text)

            print(f"  ✓ Extracted text: {len(text_parts)} text blocks from {total_slides} slides")
            return {"type": "text", "blocks": len(text_parts)}

    except Exception as e:
        print(f"  ✗ FAILED: {e}")
        import traceback
        traceback.print_exc()
        return {"type": "error", "error": str(e)}


def test_pdf_parsing(pdf_path):
    """Test PDF parsing logic (mimics parser.rs)"""
    print(f"\n--- Testing: {pdf_path.name} ---")

    try:
        import pdfplumber

        pdf = pdfplumber.open(pdf_path)
        pages = pdf.pages
        total_pages = len(pages)

        # Try to extract tables
        all_rows = []
        header = None

        for page in pages:
            tables = page.extract_tables()
            for table in tables:
                for i, row in enumerate(table):
                    if header is None and row and any(cell for cell in row):
                        header = [str(c).strip() if c else f"col_{i}" for c in row]
                    elif header is not None:
                        all_rows.append(row)

        if header and all_rows:
            print(f"  ✓ Extracted table: {len(all_rows)} rows, {len(header)} columns from {total_pages} pages")
            print(f"  Header: {header}")
            return {"type": "table", "rows": len(all_rows), "cols": len(header)}
        else:
            # Extract text
            text_parts = []
            for page in pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text.strip())

            pdf.close()
            full_text = "\n\n".join(text_parts)
            print(f"  ✓ Extracted text: {len(full_text)} characters from {total_pages} pages")
            return {"type": "text", "length": len(full_text)}

    except Exception as e:
        print(f"  ✗ FAILED: {e}")
        import traceback
        traceback.print_exc()
        return {"type": "error", "error": str(e)}


def main():
    print("="*60)
    print("Creating test files...")
    print("="*60)

    test_files = {
        "excel": [
            create_excel_merged_cells(),
            create_excel_multiple_sheets(),
            create_excel_empty_rows(),
            create_excel_large(),
        ],
        "ppt": [
            create_ppt_with_tables(),
            create_ppt_text_only(),
        ],
        "pdf": [
            create_pdf_with_table(),
            create_pdf_text_only(),
        ],
    }

    print("\n" + "="*60)
    print("Testing parsing logic...")
    print("="*60)

    results = []

    for fmt, paths in test_files.items():
        for path in paths:
            if fmt == "excel":
                result = test_excel_parsing(path)
            elif fmt == "ppt":
                result = test_ppt_parsing(path)
            elif fmt == "pdf":
                result = test_pdf_parsing(path)
            else:
                result = {"type": "unknown"}

            results.append({"file": path.name, "format": fmt, "result": result})

    print("\n" + "="*60)
    print("Summary:")
    print("="*60)

    failed = [r for r in results if r["result"]["type"] == "error"]
    if failed:
        print(f"\n✗ {len(failed)} test(s) FAILED:")
        for r in failed:
            print(f"  - {r['file']}: {r['result']['error']}")
        sys.exit(1)
    else:
        print(f"\n✓ All {len(results)} tests PASSED")
        for r in results:
            print(f"  - {r['file']} ({r['format']}): {r['result']['type']}")


if __name__ == "__main__":
    main()
