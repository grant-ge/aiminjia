#!/usr/bin/env python3
"""
End-to-end test for document loading pipeline.
Simulates the full chain: parse_file → load_file → preamble injection.

Tests ALL document types: Word, Excel, PPT, PDF, CSV, JSON, Text
Tests ALL edge cases: merged cells, multiple tables, empty files, text-only, etc.
"""

import json
import os
import sys
import traceback
from pathlib import Path

# Check dependencies
missing = []
for lib, pkg in [
    ('docx', 'python-docx'),
    ('openpyxl', 'openpyxl'),
    ('pptx', 'python-pptx'),
    ('pdfplumber', 'pdfplumber'),
    ('reportlab', 'reportlab'),
]:
    try:
        __import__(lib)
    except ImportError:
        missing.append(pkg)

if missing:
    print(f"ERROR: Missing: {', '.join(missing)}")
    sys.exit(1)

import pandas as pd
from docx import Document
from openpyxl import Workbook, load_workbook
from pptx import Presentation
from pptx.util import Inches

TEST_DIR = Path(__file__).parent / "fixtures" / "e2e"
TEST_DIR.mkdir(parents=True, exist_ok=True)

RESULTS = []
FAILED = []

# ============================================================
# Shared helpers — simulate Rust code behavior
# ============================================================

def _smart_read_data(path, sheet_name=None, nrows=None):
    """Simulates the _smart_read_data function used in preamble injection.
    This is what the actual Python runtime uses in production."""
    ext = Path(path).suffix.lower()
    if ext in ('.xlsx', '.xls'):
        wb = load_workbook(path, read_only=True, data_only=True)
        ws = wb[sheet_name] if sheet_name else wb.active
        rows = list(ws.iter_rows(max_row=(nrows or 10001), values_only=True))
        wb.close()
        if rows:
            header = [str(c).strip() if c is not None else f'col_{i}' for i, c in enumerate(rows[0])]
            return pd.DataFrame(rows[1:], columns=header)
        return pd.DataFrame()
    elif ext == '.csv':
        kwargs = {}
        if nrows:
            kwargs['nrows'] = nrows
        for enc in ['utf-8', 'utf-8-sig', 'gbk', 'latin-1']:
            try:
                return pd.read_csv(path, encoding=enc, **kwargs)
            except (UnicodeDecodeError, UnicodeError):
                continue
        return pd.read_csv(path, encoding='latin-1', errors='replace', **kwargs)
    elif ext == '.json':
        try:
            return pd.read_json(path)
        except ValueError:
            return pd.read_json(path, lines=True)
    elif ext == '.parquet':
        return pd.read_parquet(path)
    else:
        raise ValueError(f"_smart_read_data does not support '{ext}' files")


def simulate_parse_word(file_path):
    """Simulates parser.rs parse_word()"""
    doc = Document(file_path)

    # Strategy: extract the largest table (by row count)
    best_table = None
    best_row_count = 0
    for table in doc.tables:
        if len(table.rows) > best_row_count:
            best_row_count = len(table.rows)
            best_table = table

    all_rows = []
    header = None
    if best_table is not None:
        for i, row in enumerate(best_table.rows):
            seen = set()
            unique_cells = []
            for cell in row.cells:
                cell_id = id(cell)
                if cell_id not in seen:
                    seen.add(cell_id)
                    unique_cells.append(cell.text.strip())

            # Skip rows with only 1 unique cell (fully merged rows)
            if len(unique_cells) == 1:
                continue

            if header is None and any(unique_cells):
                header = unique_cells
            elif header is not None:
                n = len(header)
                if len(unique_cells) < n:
                    unique_cells.extend([''] * (n - len(unique_cells)))
                elif len(unique_cells) > n:
                    unique_cells = unique_cells[:n]
                all_rows.append(unique_cells)

    if header and all_rows:
        return {
            "type": "table",
            "columnNames": header,
            "rowCount": len(all_rows),
        }
    else:
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return {
            "type": "text",
            "paragraphCount": len(paragraphs),
            "textLength": sum(len(p) for p in paragraphs),
        }


def simulate_parse_excel(file_path):
    """Simulates parser.rs parse_file() for Excel"""
    wb = load_workbook(file_path, read_only=True, data_only=True)
    ws = wb.active
    rows = list(ws.iter_rows(max_row=10001, values_only=True))
    wb.close()
    if rows:
        header = [str(c).strip() if c is not None else f'col_{i}' for i, c in enumerate(rows[0])]
        return {"type": "table", "columnNames": header, "rowCount": max(0, ws.max_row - 1)}
    return {"type": "empty"}


def simulate_parse_ppt(file_path):
    """Simulates parser.rs parse_ppt()"""
    prs = Presentation(file_path)
    all_rows = []
    header = None
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                for i, row in enumerate(table.rows):
                    cells = [cell.text.strip() for cell in row.cells]
                    if header is None and any(cells):
                        header = cells
                    elif header is not None:
                        all_rows.append(cells)

    if header and all_rows:
        return {"type": "table", "columnNames": header, "rowCount": len(all_rows)}
    else:
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            text_parts.append(para.text.strip())
        return {"type": "text", "textBlocks": len(text_parts)}


def simulate_parse_pdf(file_path):
    """Simulates parser.rs parse_pdf()"""
    import pdfplumber
    pdf = pdfplumber.open(file_path)
    all_rows = []
    header = None
    for page in pdf.pages:
        tables = page.extract_tables()
        for table in tables:
            for i, row in enumerate(table):
                if header is None and row and any(cell for cell in row):
                    header = [str(c).strip() if c else f"col_{i}" for c in row]
                elif header is not None:
                    all_rows.append(row)

    if header and all_rows:
        pdf.close()
        return {"type": "table", "columnNames": header, "rowCount": len(all_rows)}
    else:
        text_parts = []
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                text_parts.append(text.strip())
        pdf.close()
        return {"type": "text", "textLength": sum(len(t) for t in text_parts)}


def detect_format(file_path):
    ext = Path(file_path).suffix.lower()
    mapping = {
        '.csv': 'csv', '.tsv': 'csv',
        '.xlsx': 'excel', '.xls': 'excel',
        '.json': 'json', '.jsonl': 'json',
        '.parquet': 'parquet',
        '.pdf': 'pdf',
        '.docx': 'word', '.doc': 'word',
        '.pptx': 'ppt', '.ppt': 'ppt',
        '.html': 'html', '.htm': 'html',
        '.txt': 'text', '.log': 'text',
    }
    return mapping.get(ext, 'unknown')


def determine_loaded_as(format_type, parse_result):
    """Simulates the Rust logic in file_load.rs (WITH the v0.3.19 fix)"""
    # v0.3.19 fix: Word/PPT/PDF always text
    if format_type in ('word', 'ppt', 'pdf'):
        return 'text'

    # For auto types (html)
    if format_type in ('html',):
        if parse_result.get('columnNames'):
            return 'dataframe'
        else:
            return 'text'

    # Dataframe types
    if format_type in ('csv', 'excel', 'json', 'parquet'):
        return 'dataframe'

    return 'text'


def simulate_preamble_injection(loaded_as, file_path, format_type):
    """Simulate what build_loaded_files_preamble does and verify it works"""
    if loaded_as == 'dataframe':
        # Would call _smart_read_data(path)
        try:
            df = _smart_read_data(str(file_path))
            return {"success": True, "var": "_df", "shape": list(df.shape)}
        except Exception as e:
            return {"success": False, "var": "_df", "error": str(e)}
    elif loaded_as == 'text':
        # Would call open(path, 'r', encoding='utf-8')
        # For Word/PPT, the effective_path should be a .txt file (masked or extracted)
        # For now, simulate text extraction
        if format_type == 'word':
            try:
                doc = Document(file_path)
                text = '\n'.join(p.text for p in doc.paragraphs if p.text.strip())
                return {"success": True, "var": "_text", "length": len(text)}
            except Exception as e:
                return {"success": False, "var": "_text", "error": str(e)}
        elif format_type == 'ppt':
            try:
                prs = Presentation(file_path)
                texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            texts.append(shape.text_frame.text)
                text = '\n'.join(texts)
                return {"success": True, "var": "_text", "length": len(text)}
            except Exception as e:
                return {"success": False, "var": "_text", "error": str(e)}
        elif format_type == 'pdf':
            try:
                import pdfplumber
                pdf = pdfplumber.open(file_path)
                text = '\n'.join(p.extract_text() or '' for p in pdf.pages)
                pdf.close()
                return {"success": True, "var": "_text", "length": len(text)}
            except Exception as e:
                return {"success": False, "var": "_text", "error": str(e)}
        else:
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    text = f.read()
                return {"success": True, "var": "_text", "length": len(text)}
            except Exception as e:
                return {"success": False, "var": "_text", "error": str(e)}
    else:
        return {"success": False, "error": f"Unknown loaded_as: {loaded_as}"}


def run_test(name, file_path, expected_format, expected_loaded_as):
    """Run full pipeline test for a document"""
    print(f"\n{'='*60}")
    print(f"TEST: {name}")
    print(f"  File: {file_path.name}")
    print(f"{'='*60}")

    try:
        # Step 1: Detect format
        fmt = detect_format(file_path)
        assert fmt == expected_format, f"Format mismatch: got '{fmt}', expected '{expected_format}'"
        print(f"  1. Format detection: {fmt} ✓")

        # Step 2: Parse
        if fmt == 'word':
            parse_result = simulate_parse_word(file_path)
        elif fmt == 'excel':
            parse_result = simulate_parse_excel(file_path)
        elif fmt == 'ppt':
            parse_result = simulate_parse_ppt(file_path)
        elif fmt == 'pdf':
            parse_result = simulate_parse_pdf(file_path)
        elif fmt == 'csv':
            parse_result = {"type": "table", "columnNames": ["dummy"]}
        elif fmt == 'json':
            parse_result = {"type": "table", "columnNames": ["dummy"]}
        elif fmt == 'text':
            parse_result = {"type": "text"}
        else:
            parse_result = {}

        print(f"  2. Parse result: type={parse_result.get('type', 'N/A')} ✓")

        # Step 3: Determine loaded_as (the critical fix in v0.3.19)
        loaded_as = determine_loaded_as(fmt, parse_result)
        assert loaded_as == expected_loaded_as, \
            f"loadedAs mismatch: got '{loaded_as}', expected '{expected_loaded_as}'"
        print(f"  3. loadedAs: {loaded_as} ✓")

        # Step 4: Simulate preamble injection
        preamble_result = simulate_preamble_injection(loaded_as, file_path, fmt)
        assert preamble_result["success"], \
            f"Preamble injection failed: {preamble_result.get('error', 'unknown')}"
        print(f"  4. Preamble injection: {preamble_result['var']} ✓")

        # Step 5: Verify _smart_read_data does NOT get called on Word/PPT
        if fmt in ('word', 'ppt'):
            try:
                _smart_read_data(str(file_path))
                print(f"  5. [WARN] _smart_read_data didn't crash (unexpected)")
            except Exception as e:
                print(f"  5. _smart_read_data correctly rejects {fmt}: {type(e).__name__} ✓")

        print(f"\n  ✓ PASSED: {name}")
        RESULTS.append({"name": name, "status": "PASSED"})

    except Exception as e:
        print(f"\n  ✗ FAILED: {name}")
        print(f"    Error: {e}")
        traceback.print_exc()
        RESULTS.append({"name": name, "status": "FAILED", "error": str(e)})
        FAILED.append(name)


# ============================================================
# Create test fixtures
# ============================================================

def create_fixtures():
    print("Creating test fixtures...\n")
    fixtures = {}

    # --- Word ---
    # 1. Word with table (merged cells)
    doc = Document()
    doc.add_heading('Merged Cells Test', level=1)
    table = doc.add_table(rows=4, cols=5)
    table.rows[0].cells[0].merge(table.rows[0].cells[4])
    table.rows[0].cells[0].text = 'Title Row'
    for i in range(1, 4):
        for j in range(5):
            table.rows[i].cells[j].text = f'R{i}C{j}'
    p = TEST_DIR / "word_merged.docx"
    doc.save(p)
    fixtures["word_merged"] = p

    # 2. Word text only
    doc = Document()
    doc.add_heading('Text Only', level=1)
    doc.add_paragraph('Paragraph 1')
    doc.add_paragraph('Paragraph 2')
    p = TEST_DIR / "word_text.docx"
    doc.save(p)
    fixtures["word_text"] = p

    # 3. Word with table (no merge)
    doc = Document()
    table = doc.add_table(rows=3, cols=3)
    for i, h in enumerate(['Name', 'Age', 'Dept']):
        table.rows[0].cells[i].text = h
    table.rows[1].cells[0].text = 'Alice'
    table.rows[1].cells[1].text = '30'
    table.rows[1].cells[2].text = 'Eng'
    p = TEST_DIR / "word_table.docx"
    doc.save(p)
    fixtures["word_table"] = p

    # 4. Word with complex merge
    doc = Document()
    table = doc.add_table(rows=4, cols=4)
    for i, h in enumerate(['A', 'B', 'C', 'D']):
        table.rows[0].cells[i].text = h
    table.rows[1].cells[0].merge(table.rows[1].cells[1])
    table.rows[1].cells[0].text = 'AB'
    table.rows[1].cells[2].text = 'C1'
    table.rows[1].cells[3].text = 'D1'
    for j in range(4):
        table.rows[2].cells[j].text = f'V{j}'
    p = TEST_DIR / "word_complex.docx"
    doc.save(p)
    fixtures["word_complex"] = p

    # 5. Word empty
    doc = Document()
    p = TEST_DIR / "word_empty.docx"
    doc.save(p)
    fixtures["word_empty"] = p

    # --- Excel ---
    # 6. Excel normal
    wb = Workbook()
    ws = wb.active
    ws.append(['Name', 'Value'])
    ws.append(['A', 100])
    ws.append(['B', 200])
    p = TEST_DIR / "excel_normal.xlsx"
    wb.save(p)
    fixtures["excel_normal"] = p

    # 7. Excel merged header
    wb = Workbook()
    ws = wb.active
    ws.merge_cells('A1:C1')
    ws['A1'] = 'Merged'
    for i in range(2, 5):
        for j in range(1, 4):
            ws.cell(row=i, column=j, value=f'R{i}C{j}')
    p = TEST_DIR / "excel_merged.xlsx"
    wb.save(p)
    fixtures["excel_merged"] = p

    # --- PPT ---
    # 8. PPT with table
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    table = slide.shapes.add_table(3, 3, Inches(1), Inches(1), Inches(6), Inches(2)).table
    for i, h in enumerate(['X', 'Y', 'Z']):
        table.cell(0, i).text = h
    table.cell(1, 0).text = '1'
    p = TEST_DIR / "ppt_table.pptx"
    prs.save(p)
    fixtures["ppt_table"] = p

    # 9. PPT text only
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Title"
    slide.placeholders[1].text = "Body text"
    p = TEST_DIR / "ppt_text.pptx"
    prs.save(p)
    fixtures["ppt_text"] = p

    # --- PDF ---
    from reportlab.pdfgen import canvas
    from reportlab.platypus import Table, TableStyle
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import letter

    # 10. PDF with table
    p = TEST_DIR / "pdf_table.pdf"
    c = canvas.Canvas(str(p), pagesize=letter)
    data = [['Col1', 'Col2'], ['A', '1'], ['B', '2']]
    t = Table(data)
    t.setStyle(TableStyle([('GRID', (0,0), (-1,-1), 1, colors.black)]))
    t.wrapOn(c, 400, 200)
    t.drawOn(c, 100, 600)
    c.save()
    fixtures["pdf_table"] = p

    # 11. PDF text only
    p = TEST_DIR / "pdf_text.pdf"
    c = canvas.Canvas(str(p), pagesize=letter)
    c.drawString(100, 750, "This is text only PDF.")
    c.save()
    fixtures["pdf_text"] = p

    # --- CSV ---
    # 12. CSV normal
    p = TEST_DIR / "csv_normal.csv"
    pd.DataFrame({'Name': ['A', 'B'], 'Val': [1, 2]}).to_csv(p, index=False)
    fixtures["csv_normal"] = p

    # --- JSON ---
    # 13. JSON normal
    p = TEST_DIR / "json_normal.json"
    pd.DataFrame({'X': [1, 2], 'Y': [3, 4]}).to_json(p, orient='records')
    fixtures["json_normal"] = p

    # --- Text ---
    # 14. Plain text
    p = TEST_DIR / "plain.txt"
    p.write_text("Hello world\nLine 2\n")
    fixtures["plain_text"] = p

    print(f"Created {len(fixtures)} test fixtures\n")
    return fixtures


# ============================================================
# Main
# ============================================================

def main():
    fixtures = create_fixtures()

    # Word tests — ALL must have loaded_as="text" (v0.3.19 fix)
    run_test("Word with merged cells table", fixtures["word_merged"], "word", "text")
    run_test("Word text only", fixtures["word_text"], "word", "text")
    run_test("Word with normal table", fixtures["word_table"], "word", "text")
    run_test("Word with complex merge", fixtures["word_complex"], "word", "text")
    run_test("Word empty document", fixtures["word_empty"], "word", "text")

    # Excel tests — always dataframe
    run_test("Excel normal", fixtures["excel_normal"], "excel", "dataframe")
    run_test("Excel merged header", fixtures["excel_merged"], "excel", "dataframe")

    # PPT tests — ALL must have loaded_as="text" (v0.3.19 fix)
    run_test("PPT with table", fixtures["ppt_table"], "ppt", "text")
    run_test("PPT text only", fixtures["ppt_text"], "ppt", "text")

    # PDF tests — always text (v0.3.19 fix: PDF binary not supported by _smart_read_data)
    run_test("PDF with table", fixtures["pdf_table"], "pdf", "text")
    run_test("PDF text only", fixtures["pdf_text"], "pdf", "text")

    # CSV/JSON — always dataframe
    run_test("CSV normal", fixtures["csv_normal"], "csv", "dataframe")
    run_test("JSON normal", fixtures["json_normal"], "json", "dataframe")

    # Text — always text
    run_test("Plain text", fixtures["plain_text"], "text", "text")

    # Summary
    print("\n" + "="*60)
    print(f"SUMMARY: {len(RESULTS)} tests, {len(RESULTS) - len(FAILED)} passed, {len(FAILED)} failed")
    print("="*60)

    for r in RESULTS:
        status = "✓" if r["status"] == "PASSED" else "✗"
        print(f"  {status} {r['name']}")

    if FAILED:
        print(f"\nFAILED TESTS:")
        for name in FAILED:
            print(f"  - {name}")
        sys.exit(1)
    else:
        print(f"\n✓ ALL TESTS PASSED")


if __name__ == "__main__":
    main()
